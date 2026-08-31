import collections
import io
import base64
import gc
import cv2
import numpy as np
import streamlit as st
from pptx import Presentation
from PIL import Image

st.set_page_config(page_title="Multi-Mode Image Matcher", layout="wide", initial_sidebar_state="expanded")

st.markdown("""
    <style>
    [data-testid="stToolbar"] {display: none !important;}
    footer {visibility: hidden !important;}
    [data-testid="stSidebarCollapseButton"] {display: none !important;}
    button[title="Collapse sidebar"] {display: none !important;}
    .block-container {padding-top: 1.5rem; padding-bottom: 0rem;}
    [data-testid="stSidebar"] {padding-top: 0rem;}
    h1 {font-size: 1.8rem !important; margin-bottom: 0.5rem;}
    .stAlert {padding: 0.5rem 1rem; margin-bottom: 0.5rem;}
    .zoom-img {
        width: 100px;
        border-radius: 5px;
        cursor: pointer;
        transition: transform 0.25s ease;
    }
    .zoom-img:focus {
        position: fixed;
        top: 50%;
        left: 50%;
        transform: translate(-50%, -50%) scale(4);
        max-width: 80vw;
        max-height: 80vh;
        z-index: 99999;
        box-shadow: 0 0 20px rgba(0,0,0,0.8);
        outline: none;
    }
    </style>
""", unsafe_allow_html=True)

st.title("🖼️ Dual-Mode PPT Image Matching Tool")

def get_thumbnail_base64(img, max_size=(300, 300)):
    thumb = img.copy()
    if thumb.mode in ("RGBA", "P"):
        thumb = thumb.convert("RGB")
    thumb.thumbnail(max_size)
    buffered = io.BytesIO()
    thumb.save(buffered, format="JPEG", quality=75)
    img_str = base64.b64encode(buffered.getvalue()).decode()
    return f"data:image/jpeg;base64,{img_str}"

def pil_to_cv2(pil_img):
    if pil_img.mode in ("RGBA", "P"):
        pil_img = pil_img.convert("RGB")
    open_cv_image = np.array(pil_img)
    return cv2.cvtColor(open_cv_image, cv2.COLOR_RGB2GRAY)

def get_features(cv_img):
    orb = cv2.ORB_create(nfeatures=500)
    kp, des = orb.detectAndCompute(cv_img, None)
    return des

def compare_features(des1, des2):
    if des1 is None or des2 is None:
        return 0
    bf = cv2.BFMatcher(cv2.NORM_HAMMING, crossCheck=True)
    matches = bf.match(des1, des2)
    good_matches = [m for m in matches if m.distance < 40]
    return len(good_matches)

# PPT Parsing Function with Positional Sorting
def process_ppt_file(ppt_file):
    prs = Presentation(ppt_file)
    ppt_images = []
    
    for slide_index, slide in enumerate(prs.slides):
        picture_shapes = []
        for shape in slide.shapes:
            if shape.shape_type == 13: # Picture shape type
                picture_shapes.append((shape.top, shape.left, shape))
        
        # Positional sorting (Top-to-Bottom, Left-to-Right)
        picture_shapes.sort(key=lambda x: (x[0], x[1]))
        
        for img_count, (_, _, shape) in enumerate(picture_shapes, start=1):
            try:
                image_bytes = shape.image.blob
                with Image.open(io.BytesIO(image_bytes)) as img:
                    cv_img = pil_to_cv2(img)
                    des = get_features(cv_img)
                    b64_str = get_thumbnail_base64(img)
                    
                    ppt_images.append({
                        "slide": slide_index + 1,
                        "img_num": img_count,
                        "descriptors": des,
                        "b64_img": b64_str
                    })
            except Exception:
                continue
    return ppt_images

# --- SIDEBAR CONTROLS ---
with st.sidebar:
    st.header("⚙️ Configuration")
    
    # Mode Selector
    search_mode = st.radio(
        "Search Mode Select Karein:",
        ["🖼️ Single Image Search", "📊 PPT to PPT Search"],
        index=0
    )
    st.divider()

    # Base PPT Uploader (Always active)
    main_ppt_file = st.file_uploader("1. Badi PPT Upload (600 Images)", type=["pptx"])
    st.divider()

    # Dynamic Inputs based on Mode Selection
    if search_mode == "🖼️ Single Image Search":
        single_img_file = st.file_uploader("2. Upload Search Image", type=["jpg", "jpeg", "png", "webp"])
        target_ppt_file = None
    else:
        target_ppt_file = st.file_uploader("2. Upload Chhoti PPT (10 Slides)", type=["pptx"])
        single_img_file = None

    st.divider()
    search_clicked = st.button("🚀 Start Match Search", use_container_width=True, type="primary")

# --- MAIN LOGIC ---
if main_ppt_file is None:
    st.session_state.clear()
    st.info("👈 Pehle sidebar se Badi PPT (Main File) upload karein.")
else:
    # Main PPT indexing with Cache Memory
    if "main_ppt_images" not in st.session_state or st.session_state.get("main_ppt_name") != main_ppt_file.name:
        with st.spinner("Badi PPT scan ho rahi hai (~600 images)..."):
            main_images = process_ppt_file(main_ppt_file)
            st.session_state["main_ppt_images"] = main_images
            st.session_state["main_ppt_name"] = main_ppt_file.name
            st.success(f"✅ Main PPT Indexed: {len(main_images)} images loaded!")

    main_images = st.session_state["main_ppt_images"]

    # --- MODE 1: SINGLE IMAGE SEARCH ---
    if search_mode == "🖼️ Single Image Search":
        st.subheader("🎯 Single Image Search Results")

        if search_clicked:
            if single_img_file is None:
                st.warning("⚠️ Kripya sidebar me search ke liye Image upload karein.")
            else:
                with Image.open(single_img_file) as target_img:
                    target_cv = pil_to_cv2(target_img)
                    target_des = get_features(target_cv)
                    target_b64 = get_thumbnail_base64(target_img)

                raw_matches = []
                for main_item in main_images:
                    score = compare_features(target_des, main_item["descriptors"])
                    if score >= 35:
                        raw_matches.append((main_item, score))
                
                raw_matches.sort(key=lambda x: x[1], reverse=True)

                if raw_matches:
                    highest_score = raw_matches[0][1]
                    strict_matches = [
                        (item, score) for item, score in raw_matches 
                        if score >= max(40, highest_score * 0.85)
                    ]

                    st.success(f"**Best Match Found!** Total Matches: {len(strict_matches)}")
                    c1, c2 = st.columns([1, 5])
                    with c1:
                        st.markdown(f'<img src="{target_b64}" class="zoom-img" tabindex="0">', unsafe_allow_html=True)

                    with c2:
                        for match_item, score in strict_matches:
                            st.write(f"• **Slide {match_item['slide']}** → Image #{match_item['img_num']} *(Match Score: **{score}** features)*")
                else:
                    st.error("❌ Koi bhi **90%-100% exact match** nahi mila.")
        else:
            st.info("👉 Sidebar me Image upload karke **'🚀 Start Match Search'** dabaayein.")

    # --- MODE 2: PPT TO PPT SEARCH ---
    else:
        st.subheader("📋 PPT to PPT Batch Matching Report")

        if search_clicked:
            if target_ppt_file is None:
                st.warning("⚠️ Kripya sidebar me Chhoti PPT (10 Slides) upload karein.")
            else:
                with st.spinner("Chhoti PPT scan aur match ho rahi hai..."):
                    target_images = process_ppt_file(target_ppt_file)
                    
                    if not target_images:
                        st.error("Chhoti PPT me koi images nahi mili.")
                    else:
                        found_count = 0
                        for target_item in target_images:
                            raw_matches = []
                            for main_item in main_images:
                                score = compare_features(target_item["descriptors"], main_item["descriptors"])
                                if score >= 35:
                                    raw_matches.append((main_item, score))
                            
                            raw_matches.sort(key=lambda x: x[1], reverse=True)
                            
                            if raw_matches:
                                highest_score = raw_matches[0][1]
                                strict_matches = [
                                    (item, score) for item, score in raw_matches 
                                    if score >= max(40, highest_score * 0.85)
                                ]
                            else:
                                strict_matches = []

                            st.markdown(f"#### 🎯 Chhoti PPT -> **Slide {target_item['slide']} (Image #{target_item['img_num']})**")
                            c1, c2 = st.columns([1, 5])
                            
                            with c1:
                                st.markdown(f'<img src="{target_item["b64_img"]}" class="zoom-img" tabindex="0">', unsafe_allow_html=True)
                            
                            with c2:
                                if strict_matches:
                                    found_count += 1
                                    for main_match, score in strict_matches:
                                        st.write(f"✅ Matched in **Badi PPT** → **Slide {main_match['slide']}** (Image #{main_match['img_num']}) — *Score: **{score}** features*")
                                else:
                                    st.write("❌ Badi PPT me 90%-100% match **nahi mila**.")
                            
                            st.divider()

                        st.balloons()
                        st.success(f"🎉 Complete! Total {len(target_images)} me se **{found_count} images match ho gayi**.")
        else:
            st.info("👉 Sidebar me Chhoti PPT upload karke **'🚀 Start Match Search'** dabaayein.")
